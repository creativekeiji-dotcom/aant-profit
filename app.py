import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import io
import re
import datetime
import traceback
from pptx import Presentation
from pptx.util import Inches, Pt

# ==========================================
# 1. 설정 (여기를 대폭 보강했습니다!)
# ==========================================
st.set_page_config(page_title="AANT 경영 리포트", layout="wide")

# [핵심] 이카운트에 찍히는 실제 이름들을 전부 등록
DEFAULT_FEE_RATES = {
    # 쿠팡
    "쿠팡": 0.1188, 
    "쿠팡 주식회사": 0.1188, 
    "쿠팡그로스": 0.1188, # (로켓그로스는 프로그램 내부에서 변환됨)
    
    # 네이버
    "네이버": 0.0563, # 평균치
    "네이버파이낸셜": 0.0563,
    "스마트스토어": 0.0563,
    
    # 오픈마켓 (지마켓/옥션/11번가)
    "지마켓": 0.13, # 카테고리별 상이하지만 통상 13%
    "주식회사 지마켓": 0.13, 
    "옥션": 0.13,
    "주식회사 옥션": 0.13,
    "11번가": 0.13,
    "11번가 주식회사": 0.13,
    
    # 오늘의집
    "오늘의집": 0.22,
    "버킷플레이스": 0.22,
    "(주)버킷플레이스": 0.22,

    # 기타
    "카카오톡": 0.055,
    "알리": 0.11,
    "사업자거래": 0.0,
    "기타": 0.0
}

# ==========================================
# 2. PPT 생성 함수
# ==========================================
def create_ppt(sales, gross, fixed_cost, net, margin, fig_pie, fig_bar, top10_df):
    prs = Presentation()

    # 표지
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AANT 월간 경영 분석"
    slide.placeholders[1].text = f"기준일: {datetime.date.today().strftime('%Y-%m-%d')}\n작성: 경영지원팀"

    # 요약
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "1. 경영 실적 요약"
    tf = slide.shapes.placeholders[1].text_frame
    
    def add_line(text, size, bold=False, color=None):
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        if color: p.font.color.rgb = color
        
    add_line(f"💰 총 매출액: {int(sales):,}원", 24, True)
    add_line(f"📦 매출이익: {int(gross):,}원 (이익률 {gross/sales*100:.1f}%)", 20)
    add_line(f"💸 고정비: {int(fixed_cost):,}원", 20)
    add_line(f"🏆 순이익: {int(net):,}원 (순이익률 {margin:.1f}%)", 28, True)

    # 그래프
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "2. 채널별 성과 분석"
    try:
        img_pie = fig_pie.to_image(format="png", width=500, height=400, scale=2)
        img_bar = fig_bar.to_image(format="png", width=500, height=400, scale=2)
        slide.shapes.add_picture(io.BytesIO(img_pie), Inches(0.5), Inches(2), width=Inches(4.5))
        slide.shapes.add_picture(io.BytesIO(img_bar), Inches(5.2), Inches(2), width=Inches(4.5))
    except:
        slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1)).text = "그래프 생성 실패 (서버 설정 확인 필요)"

    # 랭킹 표
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "3. 효자 상품 TOP 10 (이익금 기준)"
    if not top10_df.empty:
        rows, cols = top10_df.shape
        table = slide.shapes.add_table(rows+1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table
        for i, col in enumerate(top10_df.columns): table.cell(0, i).text = str(col)
        for i, row in top10_df.iterrows():
            for j, val in enumerate(row):
                table.cell(i+1, j).text = f"{int(val):,}" if isinstance(val, (int, float)) else str(val)
                table.cell(i+1, j).text_frame.paragraphs[0].font.size = Pt(10)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# ==========================================
# 3. 데이터 로딩 (안전 모드 유지)
# ==========================================
def safe_date_parse(val, target_year=2026):
    try:
        val_str = str(val).strip()
        match = re.search(r'(\d{1,2})/(\d{1,2})', val_str)
        if match:
            m, d = match.groups()
            return pd.to_datetime(f"{target_year}-{m}-{d}")
        return pd.to_datetime(val_str)
    except: return None

def read_file_force(file):
    try: return pd.read_excel(file, header=None, sheet_name=None)
    except: pass
    try: file.seek(0); return {'Sheet1': pd.read_csv(file, header=None, encoding='cp949')}
    except: pass
    try: file.seek(0); return {'Sheet1': pd.read_csv(file, header=None, encoding='utf-8')}
    except: return None

def load_data(files, fee_dict):
    all_dfs = []
    for file in files:
        sheets = read_file_force(file)
        if sheets is None: continue
        for name, raw in sheets.items():
            try:
                if len(raw) < 2 or raw.shape[1] < 8: continue
                # 데이터 추출
                temp = raw.iloc[:, [0, 1, 3, 4, 5, 7]].copy()
                temp.columns = ['일자_raw', '채널', '상품명', '수량', '판매단가', '원가단가']
                
                # 유효 데이터 필터링
                temp = temp[temp['일자_raw'].astype(str).str.contains(r'\d', na=False)]
                if temp.empty: continue

                # 전처리
                temp['상품명'] = temp['상품명'].fillna("상품명없음").astype(str)
                temp['채널'] = temp['채널'].fillna("기타").astype(str).str.strip() # 공백 제거 필수
                
                # 그로스 처리
                if '그로스' in str(name) or '그로스' in file.name:
                    temp['채널'] = '쿠팡그로스'
                
                all_dfs.append(temp)
            except: continue
            
    if not all_dfs: return None
    df = pd.concat(all_dfs, ignore_index=True)
    
    # 변환 및 계산
    df['일자'] = df['일자_raw'].apply(lambda x: safe_date_parse(x))
    df = df.dropna(subset=['일자'])
    df['월'] = df['일자'].dt.strftime('%Y-%m')
    for c in ['수량', '판매단가', '원가단가']:
        df[c] = pd.to_numeric(df[c], errors='coerce').fillna(0)
    
    df['총판매금액'] = df['수량'] * df['판매단가']
    df['총원가금액'] = df['수량'] * df['원가단가']
    
    # [수수료 적용 로직] - 여기가 질문하신 부분!
    # 1. 채널명 매핑
    df['수수료율'] = df['채널'].map(fee_dict)
    
    # 2. 매핑 안 된(Unknown) 채널 확인용 (NaN이면 0으로 채움)
    df['수수료율'] = df['수수료율'].fillna(0)
    
    df['수수료금액'] = df['총판매금액'] * df['수수료율']
    df['매출총이익'] = df['총판매금액'] - df['총원가금액'] - df['수수료금액']
    return df

# ==========================================
# 4. 메인 화면
# ==========================================
st.title("📊 AANT CEO 경영 대시보드")

with st.expander("📂 데이터 파일 관리", expanded=True):
    c1, c2, c3 = st.columns(3)
    up_files = c1.file_uploader("1️⃣ 판매 파일", accept_multiple_files=True, key="f1")
    cost_file = c2.file_uploader("2️⃣ 고정비 파일", key="f2")
    fee_file = c3.file_uploader("3️⃣ 수수료 파일 (업로드 시 우선 적용)", key="f3")

current_fee_rates = DEFAULT_FEE_RATES.copy()
if fee_file:
    try:
        sheets = read_file_force(fee_file)
        if sheets:
            fdf = list(sheets.values())[0]
            current_fee_rates.update(dict(zip(fdf.iloc[:, 0], fdf.iloc[:, 1])))
    except: pass

if up_files:
    try:
        df = load_data(up_files, current_fee_rates)
        if df is not None and not df.empty:
            sales = df['총판매금액'].sum()
            gross = df['매출총이익'].sum()
            fixed_cost = 0
            if cost_file:
                try:
                    sheets = read_file_force(cost_file)
                    if sheets:
                        cdf = list(sheets.values())[0]
                        fixed_cost = cdf.select_dtypes(include=['number']).sum().sum()
                except: pass
            net = gross - fixed_cost
            margin = (net / sales * 100) if sales > 0 else 0

            st.markdown("---")
            k1, k2, k3, k4 = st.columns(4)
            k1.metric("💰 총 매출", f"{int(sales):,}원")
            k2.metric("📦 매출이익", f"{int(gross):,}원")
            k3.metric("💸 고정비", f"-{int(fixed_cost):,}원")
            k4.metric("🏆 순이익", f"{int(net):,}원", delta=f"{margin:.1f}%")
            st.markdown("---")

            t1, t2, t3 = st.tabs(["📊 리포트", "✅ 수수료 검증", "💾 다운로드 (PPT/Excel)"])
            
            # --- 그래프 ---
            ch_df = df.groupby('채널')[['총판매금액', '매출총이익']].sum().reset_index()
            ch_df['이익률'] = (ch_df['매출총이익'] / ch_df['총판매금액'] * 100).fillna(0)
            ch_df = ch_df.sort_values(by='총판매금액', ascending=False)
            
            fig_pie = px.pie(ch_df, values='총판매금액', names='채널', hole=0.4, title="매출 비중")
            fig_bar = make_subplots(specs=[[{"secondary_y": True}]])
            fig_bar.add_trace(go.Bar(x=ch_df['채널'], y=ch_df['매출총이익'], name="이익금"), secondary_y=False)
            fig_bar.add_trace(go.Scatter(x=ch_df['채널'], y=ch_df['이익률'], name="이익률(%)", line=dict(color='red')), secondary_y=True)

            pr_df = df.groupby('상품명')[['수량', '총판매금액', '매출총이익']].sum().reset_index()
            pr_df = pr_df[pr_df['상품명'] != "상품명없음"]
            top10 = pd.DataFrame()
            if not pr_df.empty:
                top10 = pr_df.sort_values(by='매출총이익', ascending=False).head(10)
                top10.index = range(1, len(top10)+1)

            with t1:
                c1, c2 = st.columns([1, 2])
                with c1: st.plotly_chart(fig_pie, use_container_width=True)
                with c2: st.plotly_chart(fig_bar, use_container_width=True)
                st.divider()
                st.subheader("TOP 10 상품")
                if not top10.empty:
                    # 표 스타일링 문제 해결된 버전
                    st.dataframe(top10.style.format({'수량':'{:,.0f}','총판매금액':'{:,.0f}','매출총이익':'{:,.0f}'}), use_container_width=True)

            # [수수료 검증 탭] - 여기서 눈으로 확인 가능합니다!
            with t2:
                st.subheader("🔍 실제 적용된 수수료율 검증")
                st.info("아래 표를 보시면, 판매처별로 몇 %가 적용되었는지 한눈에 확인할 수 있습니다.")
                
                # 채널별로 평균 수수료율을 계산해서 보여줌 (제대로 매핑됐으면 설정값과 같아야 함)
                check_df = df.groupby('채널')[['총판매금액', '수수료금액']].sum().reset_index()
                check_df['실제적용률(%)'] = (check_df['수수료금액'] / check_df['총판매금액'] * 100).round(2)
                
                # 설정된 값과 비교
                st.dataframe(check_df.style.format({'총판매금액':'{:,.0f}', '수수료금액':'{:,.0f}'}), use_container_width=True)
                
                st.markdown("---")
                st.write("**⚠️ '실제적용률'이 0%로 나오는 곳이 있다면?**")
                st.write("그 사이트 이름이 '기본 수수료 목록'에 없어서 그렇습니다. 수수료 엑셀 파일을 업로드해서 이름을 추가해주세요.")

            with t3:
                st.subheader("💾 보고서 다운로드")
                
                # Excel
                buf_ex = io.BytesIO()
                with pd.ExcelWriter(buf_ex, engine='openpyxl') as writer:
                    pd.DataFrame({'구분':['매출','이익','순이익'], '금액':[sales,gross,net]}).to_excel(writer, sheet_name='요약')
                    ch_df.to_excel(writer, sheet_name='채널', index=False)
                    if not top10.empty: top10.to_excel(writer, sheet_name='랭킹', index=False)
                    df.to_excel(writer, sheet_name='전체데이터', index=False)
                
                today = datetime.date.today().strftime("%Y%m%d")
                c_d1, c_d2 = st.columns(2)
                c_d1.download_button("📥 엑셀(Excel) 다운로드", buf_ex.getvalue(), f"AANT_Report_{today}.xlsx")

                # PPT
                ppt = create_ppt(sales, gross, fixed_cost, net, margin, fig_pie, fig_bar, top10[['상품명','수량','총판매금액','매출총이익']])
                c_d2.download_button("📥 피피티(PPT) 다운로드", ppt.getvalue(), f"AANT_Report_{today}.pptx")

        else: st.error("❌ 데이터를 읽을 수 없습니다.")
    except Exception as e:
        st.error("⚠️ 오류 발생")
        st.code(traceback.format_exc())
